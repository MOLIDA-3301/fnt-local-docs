import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from io import BytesIO
from pathlib import Path

from PIL import Image
from pypdf import PdfReader, PdfWriter


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
OFFICE_EXTENSIONS = {".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".csv", ".html", ".htm"}
MAX_BATCH_BYTES = 2 * 1024 * 1024 * 1024


def find_libreoffice() -> Path | None:
    configured = os.environ.get("FNT_LIBREOFFICE_PATH")
    candidates = [
        Path(configured) if configured else None,
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
    ]
    command = shutil.which("soffice") or shutil.which("libreoffice")
    if command:
        candidates.insert(0, Path(command))
    return next((candidate for candidate in candidates if candidate and candidate.is_file()), None)


def temporary_directory() -> tempfile.TemporaryDirectory[str]:
    configured = os.environ.get("FNT_TEMP_PATH")
    return tempfile.TemporaryDirectory(dir=configured or None)


def validate_batch_size(sources: list[Path]) -> None:
    total = 0
    for source in sources:
        if not source.is_file():
            raise ValueError(f"文件不存在：{source}")
        total += source.stat().st_size
        if total > MAX_BATCH_BYTES:
            raise ValueError("一次批量任务的总输入不能超过 2GB")


def images_to_pdf(sources: list[Path], destination: Path) -> None:
    if not sources:
        raise ValueError("至少需要一张图片")
    images: list[Image.Image] = []
    total_pixels = 0
    try:
        for source in sources:
            with Image.open(source) as opened:
                width, height = opened.size
                pixels = width * height
                if pixels > 50_000_000 or width > 16384 or height > 16384:
                    raise ValueError(f"图片超出限制：{source.name}（最大 50MP / 16384px）")
                total_pixels += pixels
                if total_pixels > 100_000_000:
                    raise ValueError("图片合并任务总计超过 100MP")
                images.append(opened.convert("RGB").copy())
        destination.parent.mkdir(parents=True, exist_ok=True)
        images[0].save(destination, "PDF", save_all=True, append_images=images[1:], resolution=150)
    finally:
        for image in images:
            image.close()


def images_to_scan_pdf(sources: list[Path], destination: Path) -> None:
    from PIL import ImageEnhance, ImageOps

    if not sources:
        raise ValueError("至少需要一张图片")
    images: list[Image.Image] = []
    total_pixels = 0
    try:
        for source in sources:
            with Image.open(source) as opened:
                width, height = opened.size
                pixels = width * height
                if pixels > 50_000_000 or width > 16384 or height > 16384:
                    raise ValueError(f"图片超出限制：{source.name}（最大 50MP / 16384px）")
                total_pixels += pixels
                if total_pixels > 100_000_000:
                    raise ValueError("图片合并任务总计超过 100MP")
                gray = ImageOps.grayscale(opened)
                enhanced = ImageEnhance.Contrast(ImageOps.autocontrast(gray, cutoff=0.2)).enhance(1.12)
                images.append(enhanced.convert("RGB").copy())
        destination.parent.mkdir(parents=True, exist_ok=True)
        images[0].save(destination, "PDF", save_all=True, append_images=images[1:], resolution=150, quality=88)
    finally:
        for image in images:
            image.close()


def _font_path() -> Path:
    windows = Path(os.environ.get("WINDIR", r"C:\Windows"))
    candidates = [
        windows / "Fonts" / "msyh.ttc",
        windows / "Fonts" / "msyh.ttf",
        windows / "Fonts" / "simhei.ttf",
        windows / "Fonts" / "arial.ttf",
    ]
    return next((candidate for candidate in candidates if candidate.is_file()), candidates[-1])


def _wrap_line(draw, text: str, font, max_width: int) -> list[str]:
    if not text:
        return [""]
    lines: list[str] = []
    current = ""
    for character in text.expandtabs(4):
        candidate = current + character
        width = draw.textbbox((0, 0), candidate, font=font)[2]
        if current and width > max_width:
            lines.append(current)
            current = character
        else:
            current = candidate
    lines.append(current)
    return lines


def text_to_pdf(source: Path, destination: Path, markdown: bool = False) -> None:
    from PIL import ImageDraw, ImageFont

    content = source.read_text(encoding="utf-8-sig", errors="replace")
    if markdown:
        content = re.sub(r"!\[([^]]*)\]\([^)]*\)", r"[图片：\1]", content)
        content = re.sub(r"\[([^]]+)\]\([^)]*\)", r"\1", content)
        content = re.sub(r"^\s{0,3}(#{1,6}|>|[-+*]\s)\s*", "", content, flags=re.MULTILINE)
        content = re.sub(r"(`{1,3}|\*\*|__|~~)", "", content)
    page_size = (1240, 1754)
    margin = 100
    font = ImageFont.truetype(str(_font_path()), 28)
    line_height = 42
    probe = Image.new("RGB", page_size, "white")
    draw = ImageDraw.Draw(probe)
    wrapped = [line for raw in content.splitlines() for line in _wrap_line(draw, raw, font, page_size[0] - margin * 2)] or [""]
    lines_per_page = max(1, (page_size[1] - margin * 2) // line_height)
    pages: list[Image.Image] = []
    try:
        for start in range(0, len(wrapped), lines_per_page):
            page = Image.new("RGB", page_size, "white")
            page_draw = ImageDraw.Draw(page)
            for index, line in enumerate(wrapped[start : start + lines_per_page]):
                page_draw.text((margin, margin + index * line_height), line, fill="#1f2937", font=font)
            pages.append(page)
        destination.parent.mkdir(parents=True, exist_ok=True)
        pages[0].save(destination, "PDF", save_all=True, append_images=pages[1:], resolution=150)
    finally:
        probe.close()
        for page in pages:
            page.close()


def office_to_pdf(source: Path, destination: Path) -> None:
    executable = find_libreoffice()
    if not executable:
        raise RuntimeError("未检测到 LibreOffice。请先安装 LibreOffice，重新打开软件后再转换 Office、CSV 或 HTML 文件。")
    destination.parent.mkdir(parents=True, exist_ok=True)
    with temporary_directory() as temporary:
        temp_dir = Path(temporary)
        output_dir = temp_dir / "output"
        profile_dir = temp_dir / "profile"
        output_dir.mkdir()
        command = [
            str(executable),
            "--headless",
            "--nologo",
            "--nodefault",
            "--nolockcheck",
            f"-env:UserInstallation={profile_dir.as_uri()}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(source),
        ]
        completed = subprocess.run(command, capture_output=True, text=True, timeout=600, creationflags=0x08000000 if os.name == "nt" else 0)
        generated = output_dir / f"{source.stem}.pdf"
        if completed.returncode != 0 or not generated.is_file():
            detail = (completed.stderr or completed.stdout).strip()
            raise RuntimeError(f"LibreOffice 转换失败：{detail or source.name}")
        shutil.copyfile(generated, destination)


def source_to_pdf(source: Path, destination: Path) -> None:
    extension = source.suffix.lower()
    if extension == ".pdf":
        shutil.copyfile(source, destination)
    elif extension in IMAGE_EXTENSIONS:
        images_to_pdf([source], destination)
    elif extension in {".txt", ".md", ".markdown"}:
        text_to_pdf(source, destination, markdown=extension != ".txt")
    elif extension in OFFICE_EXTENSIONS:
        office_to_pdf(source, destination)
    else:
        raise ValueError(f"暂不支持转成 PDF 的格式：{extension or source.name}")


def files_to_pdf(sources: list[Path], destination: Path) -> None:
    if not sources:
        raise ValueError("至少需要一个文件")
    validate_batch_size(sources)
    writer = PdfWriter()
    with temporary_directory() as temporary:
        temp_dir = Path(temporary)
        for index, source in enumerate(sources):
            converted = temp_dir / f"{index:05d}.pdf"
            try:
                source_to_pdf(source, converted)
                writer.append(str(converted))
            except Exception as error:
                raise RuntimeError(f"{source.name}：{error}") from error
        destination.parent.mkdir(parents=True, exist_ok=True)
        with destination.open("wb") as stream:
            writer.write(stream)


def merge_pdfs(sources: list[Path], destination: Path) -> None:
    if not sources:
        raise ValueError("至少需要一个 PDF")
    writer = PdfWriter()
    for source in sources:
        writer.append(str(source))
    destination.parent.mkdir(parents=True, exist_ok=True)
    with destination.open("wb") as stream:
        writer.write(stream)


def split_pdf(source: Path, destination: Path, every: int, password: str | None) -> None:
    if every < 1:
        raise ValueError("每组页数必须大于 0")
    reader = PdfReader(source)
    if reader.is_encrypted and not reader.decrypt(password or ""):
        raise ValueError("PDF 密码不正确")
    destination.parent.mkdir(parents=True, exist_ok=True)
    with temporary_directory() as temp:
        temp_dir = Path(temp)
        parts: list[Path] = []
        for start in range(0, len(reader.pages), every):
            writer = PdfWriter()
            for page in reader.pages[start : start + every]:
                writer.add_page(page)
            part = temp_dir / f"{source.stem}_{start + 1:04d}-{min(start + every, len(reader.pages)):04d}.pdf"
            with part.open("wb") as stream:
                writer.write(stream)
            parts.append(part)
        with zipfile.ZipFile(destination, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for part in parts:
                archive.write(part, part.name)


def encrypt_pdf(source: Path, destination: Path, password: str) -> None:
    if not password:
        raise ValueError("密码不能为空")
    reader = PdfReader(source)
    writer = PdfWriter()
    writer.clone_document_from_reader(reader)
    writer.encrypt(user_password=password, algorithm="AES-256")
    destination.parent.mkdir(parents=True, exist_ok=True)
    with destination.open("wb") as stream:
        writer.write(stream)


def decrypt_pdf(source: Path, destination: Path, password: str) -> None:
    reader = PdfReader(source)
    if not reader.is_encrypted:
        raise ValueError("该 PDF 未加密")
    if not reader.decrypt(password):
        raise ValueError("PDF 密码不正确")
    writer = PdfWriter()
    writer.clone_document_from_reader(reader)
    destination.parent.mkdir(parents=True, exist_ok=True)
    with destination.open("wb") as stream:
        writer.write(stream)


def _open_pdf(source: Path, password: str | None) -> PdfReader:
    reader = PdfReader(source)
    if reader.is_encrypted and not reader.decrypt(password or ""):
        raise ValueError("PDF 密码不正确")
    return reader


def _parse_page_spec(specification: str | None, page_count: int) -> list[int]:
    if not specification or not specification.strip():
        return list(range(page_count))
    indexes: list[int] = []
    for part in specification.replace("，", ",").split(","):
        token = part.strip()
        if not token:
            continue
        if "-" in token:
            start_text, end_text = token.split("-", 1)
            start, end = int(start_text), int(end_text)
            step = 1 if end >= start else -1
            indexes.extend(range(start - 1, end - 1 + step, step))
        else:
            indexes.append(int(token) - 1)
    if not indexes:
        raise ValueError("页码范围不能为空")
    invalid = [index + 1 for index in indexes if index < 0 or index >= page_count]
    if invalid:
        raise ValueError(f"页码超出范围：{invalid[0]}（文档共 {page_count} 页）")
    return indexes


def organize_pdf(source: Path, destination: Path, pages: str | None, rotate: int, password: str | None) -> None:
    if rotate not in {0, 90, 180, 270}:
        raise ValueError("旋转角度只能是 0、90、180 或 270")
    reader = _open_pdf(source, password)
    indexes = _parse_page_spec(pages, len(reader.pages))
    writer = PdfWriter()
    for index in indexes:
        page = reader.pages[index]
        if rotate:
            page.rotate(rotate)
        writer.add_page(page)
    destination.parent.mkdir(parents=True, exist_ok=True)
    with destination.open("wb") as stream:
        writer.write(stream)


def compress_pdf(source: Path, destination: Path, password: str | None) -> None:
    import pymupdf as fitz

    document = fitz.open(source)
    try:
        if document.needs_pass and not document.authenticate(password or ""):
            raise ValueError("PDF 密码不正确")
        destination.parent.mkdir(parents=True, exist_ok=True)
        document.save(destination, garbage=4, deflate=True, deflate_images=True, deflate_fonts=True, clean=True)
    finally:
        document.close()


def stamp_pdf(source: Path, destination: Path, watermark: str | None, watermark_image: str | None, page_numbers: bool, password: str | None) -> None:
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen.canvas import Canvas

    if not watermark and not watermark_image and not page_numbers:
        raise ValueError("请填写水印文字、选择水印图片或启用页码")
    image_reader = None
    image_ratio = 1.0
    if watermark_image:
        image_path = Path(watermark_image)
        if not image_path.is_file() or image_path.suffix.lower() not in IMAGE_EXTENSIONS:
            raise ValueError("水印图片不存在或格式不支持")
        image_reader = ImageReader(str(image_path))
        image_width, image_height = image_reader.getSize()
        image_ratio = image_height / max(image_width, 1)
    reader = _open_pdf(source, password)
    writer = PdfWriter()
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    except KeyError:
        pass
    for index, page in enumerate(reader.pages, 1):
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        overlay_buffer = BytesIO()
        canvas = Canvas(overlay_buffer, pagesize=(width, height))
        canvas.setFont("STSong-Light", 12)
        if watermark:
            canvas.saveState()
            canvas.setFillAlpha(0.18)
            canvas.setFillColorRGB(0.35, 0.35, 0.35)
            canvas.translate(width / 2, height / 2)
            canvas.rotate(35)
            canvas.setFont("STSong-Light", max(24, min(54, width / 10)))
            canvas.drawCentredString(0, 0, watermark)
            canvas.restoreState()
        if image_reader:
            mark_width = min(width * 0.2, 130)
            mark_height = mark_width * image_ratio
            step_x = max(mark_width * 1.8, 150)
            step_y = max(mark_height * 2.2, 130)
            y = -step_y / 2
            while y < height + step_y:
                x = -step_x / 2
                while x < width + step_x:
                    canvas.saveState()
                    canvas.setFillAlpha(0.16)
                    canvas.translate(x + mark_width / 2, y + mark_height / 2)
                    canvas.rotate(28)
                    canvas.drawImage(image_reader, -mark_width / 2, -mark_height / 2, width=mark_width, height=mark_height, preserveAspectRatio=True, mask="auto")
                    canvas.restoreState()
                    x += step_x
                y += step_y
        if page_numbers:
            canvas.setFillColorRGB(0.25, 0.25, 0.25)
            canvas.drawCentredString(width / 2, 20, f"第 {index} / {len(reader.pages)} 页")
        canvas.save()
        overlay_buffer.seek(0)
        overlay = PdfReader(overlay_buffer).pages[0]
        page.merge_page(overlay)
        writer.add_page(page)
    destination.parent.mkdir(parents=True, exist_ok=True)
    with destination.open("wb") as stream:
        writer.write(stream)


def pdf_to_text(source: Path, destination: Path, output_format: str, password: str | None) -> None:
    reader = PdfReader(source)
    if reader.is_encrypted and not reader.decrypt(password or ""):
        raise ValueError("PDF 密码不正确")
    pages = [(page.extract_text() or "").strip() for page in reader.pages]
    if not any(pages):
        raise ValueError("未提取到电子文字：该文档可能是扫描件，请使用 OCR")
    if output_format == "markdown":
        content = "\n\n".join(f"## 第 {index} 页\n\n{text}" for index, text in enumerate(pages, 1) if text)
    else:
        content = "\n\n".join(f"===== 第 {index} 页 =====\n{text}" for index, text in enumerate(pages, 1) if text)
    destination.parent.mkdir(parents=True, exist_ok=True)
    destination.write_text(content, encoding="utf-8")


def pdf_to_images(source: Path, destination: Path, dpi: int, image_format: str, password: str | None) -> None:
    import pymupdf as fitz

    if dpi < 72 or dpi > 600:
        raise ValueError("图片 DPI 必须在 72 到 600 之间")
    document = fitz.open(source)
    try:
        if document.needs_pass and not document.authenticate(password or ""):
            raise ValueError("PDF 密码不正确")
        destination.parent.mkdir(parents=True, exist_ok=True)
        scale = dpi / 72
        with temporary_directory() as temporary:
            temp_dir = Path(temporary)
            rendered: list[Path] = []
            for index, page in enumerate(document, 1):
                pixmap = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
                suffix = "jpg" if image_format == "jpeg" else "png"
                output = temp_dir / f"{source.stem}_{index:04d}.{suffix}"
                pixmap.save(output)
                rendered.append(output)
            with zipfile.ZipFile(destination, "w", compression=zipfile.ZIP_DEFLATED) as archive:
                for output in rendered:
                    archive.write(output, output.name)
    finally:
        document.close()


def _ocr_pages(source: Path, password: str | None):
    import cv2
    import numpy as np
    import pymupdf as fitz
    from rapidocr import RapidOCR

    engine = RapidOCR()
    extension = source.suffix.lower()
    pages: list[tuple[object, tuple[float, float], list[dict[str, object]]]] = []
    if extension == ".pdf":
        document = fitz.open(source)
        if document.needs_pass and not document.authenticate(password or ""):
            document.close()
            raise ValueError("PDF 密码不正确")
        rendered = []
        try:
            for page in document:
                pixmap = page.get_pixmap(matrix=fitz.Matrix(200 / 72, 200 / 72), alpha=False)
                image = np.frombuffer(pixmap.samples, dtype=np.uint8).reshape(pixmap.height, pixmap.width, pixmap.n).copy()
                image = cv2.cvtColor(image, cv2.COLOR_RGB2BGR)
                rendered.append((image, (page.rect.width, page.rect.height)))
        finally:
            document.close()
    elif extension in IMAGE_EXTENSIONS:
        image = cv2.imread(str(source), cv2.IMREAD_COLOR)
        if image is None:
            raise ValueError(f"无法读取图片：{source.name}")
        rendered = [(image, (float(image.shape[1]), float(image.shape[0])))]
    else:
        raise ValueError("OCR 仅支持图片或 PDF")

    for image, page_size in rendered:
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY) if len(image.shape) == 3 else image
        normalized = cv2.normalize(gray, None, 0, 255, cv2.NORM_MINMAX)
        result = engine(normalized)
        records: list[dict[str, object]] = []
        if result and result.txts:
            for box, text, score in zip(result.boxes, result.txts, result.scores):
                records.append({"box": box.tolist(), "text": text, "score": float(score)})
        pages.append((image, page_size, records))
    return pages


def ocr_to_output(source: Path, destination: Path, output_format: str, password: str | None, min_confidence: float) -> None:
    import cv2
    import pymupdf as fitz

    pages = _ocr_pages(source, password)
    destination.parent.mkdir(parents=True, exist_ok=True)

    if output_format in {"text", "markdown"}:
        page_blocks: list[str] = []
        for page_number, (_, _, records) in enumerate(pages, 1):
            lines = []
            for record in records:
                text = str(record["text"])
                score = float(record["score"])
                lines.append(text if score >= min_confidence else f"[低置信度 {score:.0%}] {text}")
            heading = f"## 第 {page_number} 页" if output_format == "markdown" else f"===== 第 {page_number} 页 ====="
            page_blocks.append(f"{heading}\n\n" + "\n".join(lines))
        destination.write_text("\n\n".join(page_blocks), encoding="utf-8")
        return

    if output_format == "docx":
        from docx import Document
        from docx.enum.text import WD_COLOR_INDEX

        document = Document()
        for page_number, (_, _, records) in enumerate(pages, 1):
            document.add_heading(f"第 {page_number} 页", level=1)
            for record in records:
                paragraph = document.add_paragraph()
                run = paragraph.add_run(str(record["text"]))
                score = float(record["score"])
                if score < min_confidence:
                    run.font.highlight_color = WD_COLOR_INDEX.YELLOW
                    run.add_text(f"  [低置信度 {score:.0%}]")
        document.save(destination)
        return

    if output_format == "searchable-pdf":
        output = fitz.open()
        font_path = _font_path()
        for image, page_size, records in pages:
            page_width, page_height = page_size
            page = output.new_page(width=page_width, height=page_height)
            ok, encoded = cv2.imencode(".jpg", image, [int(cv2.IMWRITE_JPEG_QUALITY), 92])
            if not ok:
                raise RuntimeError("无法编码 OCR 页面")
            page.insert_image(page.rect, stream=encoded.tobytes())
            font_name = "fnt-cjk"
            page.insert_font(fontname=font_name, fontfile=str(font_path))
            image_height, image_width = image.shape[:2]
            scale_x = page_width / image_width
            scale_y = page_height / image_height
            for record in records:
                box = record["box"]
                x0 = min(point[0] for point in box) * scale_x
                y0 = min(point[1] for point in box) * scale_y
                x1 = max(point[0] for point in box) * scale_x
                y1 = max(point[1] for point in box) * scale_y
                rect = fitz.Rect(x0, y0, max(x0 + 2, x1), max(y0 + 2, y1))
                page.insert_text((rect.x0, max(rect.y0 + 4, rect.y1 - 1)), str(record["text"]), fontname=font_name, fontsize=max(4, rect.height * 0.72), render_mode=3)
        output.save(destination, garbage=4, deflate=True)
        output.close()
        return

    raise ValueError(f"不支持的 OCR 输出格式：{output_format}")


def pdf_to_word(source: Path, destination: Path, password: str | None) -> None:
    from pdf2docx import Converter

    reader = PdfReader(source)
    if reader.is_encrypted and not reader.decrypt(password or ""):
        raise ValueError("PDF 密码不正确")
    if not any((page.extract_text() or "").strip() for page in reader.pages):
        ocr_to_output(source, destination, "docx", password, 0.80)
        return
    destination.parent.mkdir(parents=True, exist_ok=True)
    converter = Converter(str(source), password=password or None)
    try:
        converter.convert(str(destination))
    finally:
        converter.close()


def _ocr_records_to_rows(records: list[dict[str, object]]) -> list[list[tuple[str, float]]]:
    if not records:
        return []
    ordered = sorted(records, key=lambda record: (sum(point[1] for point in record["box"]) / 4, min(point[0] for point in record["box"])))
    rows: list[list[dict[str, object]]] = []
    for record in ordered:
        center_y = sum(point[1] for point in record["box"]) / 4
        height = max(point[1] for point in record["box"]) - min(point[1] for point in record["box"])
        if not rows:
            rows.append([record])
            continue
        previous_y = sum(point[1] for point in rows[-1][0]["box"]) / 4
        if abs(center_y - previous_y) <= max(8, height * 0.7):
            rows[-1].append(record)
        else:
            rows.append([record])
    return [[(str(record["text"]), float(record["score"])) for record in sorted(row, key=lambda item: min(point[0] for point in item["box"]))] for row in rows]


def pdf_to_excel(source: Path, destination: Path, password: str | None, min_confidence: float) -> None:
    import pdfplumber
    from openpyxl import Workbook
    from openpyxl.comments import Comment
    from openpyxl.styles import Alignment, Font, PatternFill

    tables: list[tuple[int, list[list[tuple[str, float]]], str]] = []
    raw_rows: list[list[object]] = [["页码", "x0", "top", "x1", "bottom", "文字", "置信度", "来源"]]
    scanned_pages: list[int] = []
    with pdfplumber.open(source, password=password or None) as document:
        for page_number, page in enumerate(document.pages, 1):
            words = page.extract_words(use_text_flow=True, keep_blank_chars=False)
            for word in words:
                raw_rows.append([page_number, word.get("x0"), word.get("top"), word.get("x1"), word.get("bottom"), word.get("text", ""), 1.0, "电子文字"])
            if not words:
                scanned_pages.append(page_number)
                continue
            found: list[list[list[str | None]]] = []
            for settings in ({}, {"vertical_strategy": "text", "horizontal_strategy": "text", "min_words_vertical": 2, "min_words_horizontal": 1}):
                for table in page.extract_tables(table_settings=settings):
                    normalized = [[(cell or "").strip() for cell in row] for row in table if any((cell or "").strip() for cell in row)]
                    if normalized and normalized not in found:
                        found.append(normalized)
            for table in found:
                rows = [[(cell, 1.0) for cell in row] for row in table]
                tables.append((page_number, rows, "电子表格"))

    if scanned_pages:
        ocr_pages = _ocr_pages(source, password)
        for page_number in scanned_pages:
            records = ocr_pages[page_number - 1][2]
            for record in records:
                box = record["box"]
                raw_rows.append([page_number, min(point[0] for point in box), min(point[1] for point in box), max(point[0] for point in box), max(point[1] for point in box), record["text"], record["score"], "OCR"])
            rows = _ocr_records_to_rows(records)
            if rows:
                tables.append((page_number, rows, "OCR 无框表格回退"))

    workbook = Workbook()
    workbook.remove(workbook.active)
    header_fill = PatternFill("solid", fgColor="E8E6FF")
    low_fill = PatternFill("solid", fgColor="FFF1A8")
    for table_index, (page_number, rows, source_kind) in enumerate(tables, 1):
        sheet = workbook.create_sheet(f"表{table_index}_页{page_number}"[:31])
        sheet.freeze_panes = "A2"
        sheet.sheet_properties.pageSetUpPr.fitToPage = True
        for row_index, row in enumerate(rows, 1):
            for column_index, (value, score) in enumerate(row, 1):
                cell = sheet.cell(row=row_index, column=column_index, value=value)
                cell.alignment = Alignment(vertical="top", wrap_text=True)
                if row_index == 1:
                    cell.font = Font(bold=True)
                    cell.fill = header_fill
                if score < min_confidence:
                    cell.fill = low_fill
                    cell.comment = Comment(f"OCR 置信度 {score:.1%}，请人工核对。来源：{source_kind}", "FNT")
        for column in sheet.columns:
            maximum = min(48, max(10, max(len(str(cell.value or "")) for cell in column) + 2))
            sheet.column_dimensions[column[0].column_letter].width = maximum
        sheet["A1"].comment = Comment(f"来源页：{page_number}；提取方式：{source_kind}。空白单元格可能对应合并区域，请结合 Raw 表核对。", "FNT")

    raw = workbook.create_sheet("Raw")
    for row in raw_rows:
        raw.append(row)
    raw.freeze_panes = "A2"
    for cell in raw[1]:
        cell.font = Font(bold=True)
        cell.fill = header_fill
    widths = [8, 12, 12, 12, 12, 42, 12, 14]
    for index, width in enumerate(widths, 1):
        raw.column_dimensions[raw.cell(1, index).column_letter].width = width
    if not tables:
        notice = workbook.create_sheet("说明", 0)
        notice["A1"] = "未识别到结构化表格，请查看 Raw 工作表中的坐标文字。"
        notice.column_dimensions["A"].width = 72
    destination.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(destination)


def pdf_to_ppt(source: Path, destination: Path, password: str | None, dpi: int) -> None:
    import pymupdf as fitz
    from pptx import Presentation
    from pptx.util import Inches

    if dpi < 96 or dpi > 300:
        raise ValueError("PPT 渲染 DPI 必须在 96 到 300 之间")
    pdf = fitz.open(source)
    if pdf.needs_pass and not pdf.authenticate(password or ""):
        pdf.close()
        raise ValueError("PDF 密码不正确")
    if len(pdf) == 0:
        pdf.close()
        raise ValueError("PDF 没有页面")
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    first_ratio = pdf[0].rect.height / pdf[0].rect.width
    presentation.slide_height = int(presentation.slide_width * first_ratio)
    while presentation.slides:
        relationship_id = presentation.slides._sldIdLst[-1].rId
        presentation.part.drop_rel(relationship_id)
        del presentation.slides._sldIdLst[-1]
    with temporary_directory() as temporary:
        temp_dir = Path(temporary)
        for index, page in enumerate(pdf, 1):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72), alpha=False)
            image_path = temp_dir / f"page-{index:04d}.png"
            pixmap.save(image_path)
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            image_ratio = page.rect.width / page.rect.height
            slide_ratio = presentation.slide_width / presentation.slide_height
            if image_ratio > slide_ratio:
                width = presentation.slide_width
                height = int(width / image_ratio)
                left = 0
                top = int((presentation.slide_height - height) / 2)
            else:
                height = presentation.slide_height
                width = int(height * image_ratio)
                top = 0
                left = int((presentation.slide_width - width) / 2)
            slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    pdf.close()
    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(destination)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser()
    subparsers = parser.add_subparsers(dest="command", required=True)
    images = subparsers.add_parser("images-to-pdf")
    images.add_argument("--destination", required=True)
    images.add_argument("--source", action="append", required=True)
    scan_images = subparsers.add_parser("images-to-scan-pdf")
    scan_images.add_argument("--destination", required=True)
    scan_images.add_argument("--source", action="append", required=True)
    merge = subparsers.add_parser("merge-pdf")
    merge.add_argument("--destination", required=True)
    merge.add_argument("--source", action="append", required=True)
    split = subparsers.add_parser("split-pdf")
    split.add_argument("--source", required=True)
    split.add_argument("--destination", required=True)
    split.add_argument("--every", type=int, default=1)
    split.add_argument("--password")
    encrypt = subparsers.add_parser("encrypt-pdf")
    encrypt.add_argument("--source", required=True)
    encrypt.add_argument("--destination", required=True)
    encrypt.add_argument("--password", required=True)
    decrypt = subparsers.add_parser("decrypt-pdf")
    decrypt.add_argument("--source", required=True)
    decrypt.add_argument("--destination", required=True)
    decrypt.add_argument("--password", required=True)
    organize = subparsers.add_parser("organize-pdf")
    organize.add_argument("--source", required=True)
    organize.add_argument("--destination", required=True)
    organize.add_argument("--pages")
    organize.add_argument("--rotate", type=int, default=0)
    organize.add_argument("--password")
    compress = subparsers.add_parser("compress-pdf")
    compress.add_argument("--source", required=True)
    compress.add_argument("--destination", required=True)
    compress.add_argument("--password")
    stamp = subparsers.add_parser("stamp-pdf")
    stamp.add_argument("--source", required=True)
    stamp.add_argument("--destination", required=True)
    stamp.add_argument("--watermark")
    stamp.add_argument("--watermark-image")
    stamp.add_argument("--page-numbers", action="store_true")
    stamp.add_argument("--password")
    text = subparsers.add_parser("pdf-to-text")
    text.add_argument("--source", required=True)
    text.add_argument("--destination", required=True)
    text.add_argument("--format", choices=["text", "markdown"], default="text")
    text.add_argument("--password")
    files = subparsers.add_parser("files-to-pdf")
    files.add_argument("--destination", required=True)
    files.add_argument("--source", action="append", required=True)
    subparsers.add_parser("engine-status")
    images_export = subparsers.add_parser("pdf-to-images")
    images_export.add_argument("--source", required=True)
    images_export.add_argument("--destination", required=True)
    images_export.add_argument("--dpi", type=int, default=150)
    images_export.add_argument("--format", choices=["png", "jpeg"], default="png")
    images_export.add_argument("--password")
    word = subparsers.add_parser("pdf-to-word")
    word.add_argument("--source", required=True)
    word.add_argument("--destination", required=True)
    word.add_argument("--password")
    ocr = subparsers.add_parser("ocr-document")
    ocr.add_argument("--source", required=True)
    ocr.add_argument("--destination", required=True)
    ocr.add_argument("--format", choices=["text", "markdown", "docx", "searchable-pdf"], required=True)
    ocr.add_argument("--password")
    ocr.add_argument("--min-confidence", type=float, default=0.80)
    excel = subparsers.add_parser("pdf-to-excel")
    excel.add_argument("--source", required=True)
    excel.add_argument("--destination", required=True)
    excel.add_argument("--password")
    excel.add_argument("--min-confidence", type=float, default=0.80)
    ppt = subparsers.add_parser("pdf-to-ppt")
    ppt.add_argument("--source", required=True)
    ppt.add_argument("--destination", required=True)
    ppt.add_argument("--password")
    ppt.add_argument("--dpi", type=int, default=150)
    return parser


def main() -> None:
    args = build_parser().parse_args()
    result: dict[str, object] = {"ok": True}
    if args.command == "images-to-pdf":
        images_to_pdf([Path(value) for value in args.source], Path(args.destination))
    elif args.command == "images-to-scan-pdf":
        images_to_scan_pdf([Path(value) for value in args.source], Path(args.destination))
    elif args.command == "merge-pdf":
        merge_pdfs([Path(value) for value in args.source], Path(args.destination))
    elif args.command == "split-pdf":
        split_pdf(Path(args.source), Path(args.destination), args.every, args.password)
    elif args.command == "encrypt-pdf":
        encrypt_pdf(Path(args.source), Path(args.destination), args.password)
    elif args.command == "decrypt-pdf":
        decrypt_pdf(Path(args.source), Path(args.destination), args.password)
    elif args.command == "organize-pdf":
        organize_pdf(Path(args.source), Path(args.destination), args.pages, args.rotate, args.password)
    elif args.command == "compress-pdf":
        compress_pdf(Path(args.source), Path(args.destination), args.password)
    elif args.command == "stamp-pdf":
        stamp_pdf(Path(args.source), Path(args.destination), args.watermark, args.watermark_image, args.page_numbers, args.password)
    elif args.command == "pdf-to-text":
        pdf_to_text(Path(args.source), Path(args.destination), args.format, args.password)
    elif args.command == "files-to-pdf":
        files_to_pdf([Path(value) for value in args.source], Path(args.destination))
    elif args.command == "engine-status":
        libreoffice = find_libreoffice()
        result["libreoffice"] = str(libreoffice) if libreoffice else None
    elif args.command == "pdf-to-images":
        pdf_to_images(Path(args.source), Path(args.destination), args.dpi, args.format, args.password)
    elif args.command == "pdf-to-word":
        pdf_to_word(Path(args.source), Path(args.destination), args.password)
    elif args.command == "ocr-document":
        ocr_to_output(Path(args.source), Path(args.destination), args.format, args.password, args.min_confidence)
    elif args.command == "pdf-to-excel":
        pdf_to_excel(Path(args.source), Path(args.destination), args.password, args.min_confidence)
    elif args.command == "pdf-to-ppt":
        pdf_to_ppt(Path(args.source), Path(args.destination), args.password, args.dpi)
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    try:
        main()
    except Exception as error:
        print(f"{type(error).__name__}: {error}", file=sys.stderr)
        raise SystemExit(1)
